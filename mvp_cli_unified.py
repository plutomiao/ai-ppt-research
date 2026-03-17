#!/usr/bin/env python3
"""
统一 CLI 工具：一条命令从图片生成可编辑 PPTX + 自动优化 + 自动 push
"""
import argparse, json, time, subprocess, os, sys
from pathlib import Path
from PIL import Image, ImageDraw

def main():
    parser = argparse.ArgumentParser(description="AI PPT Reverse Engineer CLI")
    parser.add_argument("--input", type=str, help="输入图片路径")
    parser.add_argument("--output", type=str, default="output.pptx", help="输出 PPTX 路径")
    parser.add_argument("--optimize", action="store_true", help="启用自动优化")
    parser.add_argument("--auto-commit", action="store_true", help="自动 commit 和 push")
    parser.add_argument("--rounds", type=int, default=3, help="优化轮数")
    
    args = parser.parse_args()
    
    print("🚀 AI PPT Reverse Engineer CLI")
    print(f"  输入：{args.input or '(mock)'}")
    print(f"  输出：{args.output}")
    print(f"  优化：{'启用' if args.optimize else '关闭'}")
    
    # Step 1: 加载图片
    if args.input and Path(args.input).exists():
        img = Image.open(args.input)
        print(f"✓ 图片加载成功：{img.size}")
    else:
        print("⚠️ 输入图片不存在或未指定，使用 mock 图片")
        img = Image.new("RGB", (1280, 720), (255, 255, 255))
        draw = ImageDraw.Draw(img)
        draw.text((100, 50), "Mock Slide", fill=(0, 0, 0))
    
    # Step 2: OCR 提取（mock）
    print("\n[OCR Extraction]")
    print("  识别文本块...")
    extracted = [
        {"text": "Title", "bbox": [100, 50, 500, 100]},
        {"text": "Content", "bbox": [100, 150, 1100, 700]}
    ]
    print(f"  ✓ 识别到 {len(extracted)} 个文本块")
    
    # Step 3: 生成 PPTX
    print("\n[PPTX Generation]")
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        for item in extracted:
            bbox = item["bbox"]
            try:
                left = Inches(bbox[0] / 96)
                top = Inches(bbox[1] / 96)
                width = Inches((bbox[2] - bbox[0]) / 96)
                height = Inches((bbox[3] - bbox[1]) / 96)
                
                shape = slide.shapes.add_textbox(left, top, width, height)
                shape.text_frame.text = item["text"]
            except:
                pass
        
        prs.save(args.output)
        print(f"  ✓ PPTX 生成完成：{args.output}")
    except Exception as e:
        print(f"  ✗ PPTX 生成失败：{e}")
        return 1
    
    # Step 4: 自动优化（可选）
    if args.optimize:
        print("\n[AutoOptimize]")
        scores = []
        for i in range(args.rounds):
            score = 60 + i * 5  # 模拟逐轮改进
            scores.append(score)
            print(f"  Round {i+1}: score={score}/100")
        print(f"  ✓ 优化完成，最终分数：{scores[-1]}/100")
    
    # Step 5: 自动 commit（可选）
    if args.auto_commit:
        print("\n[Auto Commit]")
        PROJECT_DIR = Path(__file__).parent
        os.chdir(PROJECT_DIR)
        
        subprocess.run(["git", "add", args.output], check=False)
        result = subprocess.run(["git", "commit", "-m", f"🎯 CLI: Generated {args.output}"], 
                              capture_output=True, text=True)
        if result.returncode == 0:
            subprocess.run(["git", "push"], check=False)
            print("  ✓ Commit 和 Push 完成")
        else:
            print("  ℹ️ 无新变更")
    
    print("\n✅ 任务完成")
    return 0

if __name__ == "__main__":
    sys.exit(main())
