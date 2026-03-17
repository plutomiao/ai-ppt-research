from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from .config import BOX_FILL, BOX_OUTLINE, SLIDE_SIZE, TITLE_COLOR
from .models import SlideScene


def _px_to_inches(value: int) -> float:
    return value / 96.0


def export_scene_to_pptx(scene: SlideScene, output_path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(_px_to_inches(SLIDE_SIZE[0]))
    presentation.slide_height = Inches(_px_to_inches(SLIDE_SIZE[1]))

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(_px_to_inches(scene.box.x)),
        Inches(_px_to_inches(scene.box.y)),
        Inches(_px_to_inches(scene.box.w)),
        Inches(_px_to_inches(scene.box.h)),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*BOX_FILL)
    box.line.color.rgb = RGBColor(*BOX_OUTLINE)
    box.line.width = Pt(2.5)

    title_box = slide.shapes.add_textbox(
        Inches(_px_to_inches(scene.title.x)),
        Inches(_px_to_inches(scene.title.y)),
        Inches(_px_to_inches(430)),
        Inches(_px_to_inches(80)),
    )
    paragraph = title_box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = scene.title.text
    run.font.size = Pt(36)
    run.font.color.rgb = RGBColor(*TITLE_COLOR)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
